from pptx import Presentation
from pptx.util import Inches, Pt

def add_slide(prs, title_text, bullet_points):
    """Helper function to add a slide with a title and bullet points."""
    slide_layout = prs.slide_layouts[1] # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = title_text
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.word_wrap = True
    
    for i, point in enumerate(bullet_points):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = point
        p.level = 0
        p.font.size = Pt(18)

# Initialize presentation
prs = Presentation()

# --- Slide 1: Title Slide ---
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "NEXT-GEN AI-POWERED RECRUITMENT & ANALYTICS ECOSYSTEM"
subtitle.text = "Grounded in IEEE: 'Next-Gen Recruitment: An AI Powered Hiring Ecosystem using NLP'\n\nTeam Presentation"

# --- Slide 2: The Problem Statement ---
add_slide(prs, "The Problem in Traditional Recruitment", [
    "Manual resume screening is time-consuming and introduces severe human bias.",
    "Traditional ATS relies on binary, exact-match keyword searching.",
    "Qualified candidates are frequently rejected due to formatting variations rather than lack of underlying skills.",
    "No feedback loop: Students are rejected without understanding their specific technical skill gaps."
])

# --- Slide 3: Academic Foundation & IEEE Significance ---
add_slide(prs, "Academic Foundation (IEEE)", [
    "This architecture is strictly grounded in the IEEE publication:",
    "\"Next-Gen Recruitment: An AI Powered Hiring Ecosystem using NLP\"",
    "Core Shift: Moving from basic keyword 'Searching' to semantic 'Intelligence'.",
    "Dual Purpose: Functions as a high-efficiency recruitment tool for administrators AND a career-development mentor for students."
])

# --- Slide 4: Proposed Solution & AI Flow (In & Out) ---
add_slide(prs, "AI Mechanics: The Data Flow", [
    "Input: Unstructured Student Resume (PDF) & Target Job Description (PDF).",
    "NLP Pipeline: Utilizes spaCy for Named Entity Recognition (NER) to extract skills, experience, and academic metrics.",
    "Semantic Analysis: Sentence-Transformers (BERT) maps text to high-dimensional vectors to calculate contextual similarity.",
    "Output: A deterministic Hybrid Match Score and a precise array of Missing Technologies."
])

# --- Slide 5: Interactive Skill Gap & Mock Interviews ---
add_slide(prs, "Mentorship Features (IEEE Grounded)", [
    "Interactive Skill Gap Analysis: Outputs specific technologies a student must learn to reach a 90% compatibility threshold.",
    "Mock Interview Module: Integrates with Generative AI (Gemini/LLMs).",
    "Simulates technical interviews based strictly on the student's extracted tech stack and their identified missing skills.",
    "Fulfills the IEEE mandate of acting as a career-development mentor."
])

# --- Slide 6: Backend & Frontend Architecture ---
add_slide(prs, "System Design: Microservices", [
    "Frontend (UI): Built with HTML5, CSS3, and Bootstrap for a responsive administrator dashboard. Implements Chart.js for batch data visualization.",
    "Primary Backend: Node.js server acts as the central router and API gateway.",
    "AI Microservice: Python FastAPI server executing the heavy NLP/BERT mathematical operations asynchronously.",
    "Database: PostgreSQL stores candidate dossiers, multi-factor rankings (GPA + Skills), and historical match analytics."
])

# --- Slide 7: Conclusion ---
add_slide(prs, "Conclusion & Expected Impact", [
    "Replaces manual processing with mathematically verified semantic analysis.",
    "Provides actionable data to both the Placement Cell and the student body.",
    "Strictly adheres to the methodologies outlined in the referenced IEEE literature.",
    "Scalable architecture designed for modern cloud deployment."
])

# Save the presentation
filename = "Next_Gen_Recruitment_Presentation.pptx"
prs.save(filename)
print(f"Presentation successfully generated: {filename}")